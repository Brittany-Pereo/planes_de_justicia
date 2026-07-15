# -*- coding: utf-8 -*-
"""crear_pptx
====================================================================
Conversión a Python de R/utils_crear_pptx.R

Genera el reporte de productividad del IMSS Bienestar en PowerPoint.

Equivalencias de librerías respecto al código R original:
    - officer                    -> python-pptx (lectura/escritura del .pptx)
    - rvg::dml (tarjetas/gráficas) -> formas nativas de python-pptx
      (rectángulos, líneas freeform, conectores, texto — todo editable
      en PowerPoint, igual que rvg::dml convertía ggplot2 a DrawingML)
    - flextable                  -> tablas nativas de python-pptx
    - dplyr / tidyr              -> pandas

La función principal `crear_reporte_productividad` recibe los mismos
DataFrames que la versión de R y devuelve el objeto `Presentation`
(equivalente al objeto `pptx` que devolvía R). Para guardar:

    prs = crear_reporte_productividad(...)
    prs.save("reporte.pptx")

Columnas esperadas de los DataFrames de entrada (igual que en R):
    historicos:  fecha, consulta_total, consulta_general,
                 consulta_especialidad, procedimientos_qx, egresos
                 (columnas faltantes se tratan como 0)
    metas:       clues, meta_general_anual, meta_especialidad_anual,
                 meta_cirugia_anual
    clues_info:  clues, id, nombre, entidad
    procedimientos_personas: fecha (año), tipo_procedimiento,
                 procedimientos, personas   (o None para derivarlo)
"""

from __future__ import annotations

import math
from datetime import date, timedelta
from pathlib import Path

import numpy as np
import pandas as pd

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.oxml.ns import qn

# Meses en español (evita depender del locale del sistema en Windows)
MESES_ES = [
    "enero", "febrero", "marzo", "abril", "mayo", "junio",
    "julio", "agosto", "septiembre", "octubre", "noviembre", "diciembre",
]


# Colores ------------------------------------------------------------------
col_rojo_chillon = "#FF0000"
col_amarillo_chillon = "#FFC107"
col_verde_chillon = "#00B050"
# Color de letras y bordes
col_muted = "#6B7280"
col_borde = "#D1D5DB"
col_texto = "#111827"
# Color de gráficas
col_verde = "#1E5B4F"   # IMSS verde
col_guinda = "#611232"  # guinda
col_dorado = "#A57F2C"
col_verde_pastel = "#FFFFFF"


# ---------------------------------------------------------------------------
# Utilidades generales
# ---------------------------------------------------------------------------
def _rgb(hex_color: str) -> RGBColor:
    """Convierte '#RRGGBB' a RGBColor de python-pptx."""
    return RGBColor.from_string(hex_color.lstrip("#").upper())


def fmt_num(x) -> str:
    """scales::comma(as.integer(x)) -> '1,234,567'."""
    try:
        return f"{int(round(float(x))):,}"
    except (TypeError, ValueError):
        return "0"


def _es_vacio(x) -> bool:
    if x is None:
        return True
    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        return len(x) == 0
    return False


def _primer_valor(x):
    """Equivalente a x[1] de R (primer elemento) para escalares o vectores."""
    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        return x[0] if len(x) else np.nan
    return x


# ---------------------------------------------------------------------------
# Funciones de valuebox
# ---------------------------------------------------------------------------
def fmt_delta(x) -> dict:
    """Formatea la variación porcentual: etiqueta, color e ícono."""
    if x is None:
        return {"label": "", "col": col_muted, "icon": ""}

    if isinstance(x, (list, tuple, np.ndarray, pd.Series)):
        if len(x) == 0:
            return {"label": "", "col": col_muted, "icon": ""}
        if pd.isna(pd.Series(list(x))).all():
            return {"label": "s/d", "col": col_muted, "icon": ""}
        x = x[0]

    if pd.isna(x):
        return {"label": "s/d", "col": col_muted, "icon": ""}

    x = _primer_valor(x)

    if x > 0:
        return {"label": f"+{fmt_pct_int(x)}%", "col": col_verde, "icon": "▲ "}
    if x < 0:
        return {"label": f"{fmt_pct_int(x)}%", "col": col_guinda, "icon": "▼ "}
    return {"label": "0%", "col": col_muted, "icon": "• "}


def fmt_pct_int(x) -> str:
    """Representa el porcentaje sin decimales (como R con round(...,0))."""
    return str(int(round(float(x))))


def elige_acento(var_2025, var_2024,
                 verde=col_verde_chillon,
                 amarillo=col_amarillo_chillon,
                 rojo=col_rojo_chillon) -> str:
    """Elige el color de acento según cuántas variaciones son negativas."""
    v25_neg = (var_2025 is not None) and (not pd.isna(var_2025)) and var_2025 < 0
    v24_neg = (var_2024 is not None) and (not pd.isna(var_2024)) and var_2024 < 0

    if v25_neg and v24_neg:
        return rojo
    if v25_neg != v24_neg:  # xor
        return amarillo
    return verde


def crear_card_institucional(numero, titulo, var_vs_2025, var_vs_2024,
                             acento=col_verde,
                             size_num=18, size_titulo=11, size_delta=9) -> dict:
    """Devuelve una especificación de tarjeta (se dibuja al colocarla).

    En R esto devolvía un objeto `rvg::dml`; aquí devolvemos un dict con la
    información necesaria para dibujar la tarjeta como formas nativas dentro
    del placeholder de destino (ver `dibujar_card`).
    """
    d25 = fmt_delta(var_vs_2025)
    d24 = fmt_delta(var_vs_2024)

    mostrar_comparativos = (var_vs_2025 is not None) and (var_vs_2024 is not None)

    return {
        "numero": fmt_num(numero),
        "titulo": titulo,
        "d25": d25,
        "d24": d24,
        "acento": acento,
        "mostrar_comparativos": mostrar_comparativos,
        "size_num": size_num,
        "size_titulo": size_titulo,
        "size_delta": size_delta,
    }


def dibujar_card(slide, box, spec):
    """Dibuja la tarjeta como formas nativas dentro del rectángulo `box`.

    box = (left, top, width, height) en EMU.
    """
    L, T, W, H = box

    # --- Card base (rectángulo redondeado)
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, L, T, W, H)
    card.fill.solid()
    card.fill.fore_color.rgb = _rgb("#FFFFFF")
    card.line.color.rgb = _rgb(col_borde)
    card.line.width = Pt(1)
    card.shadow.inherit = False

    # --- Línea lateral (acento)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        L + int(0.025 * W), T + int(0.05 * H),
        int(0.02 * W), int(0.90 * H),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(spec["acento"])
    bar.line.fill.background()
    bar.shadow.inherit = False

    def _txt(top_frac, height_frac):
        tb = slide.shapes.add_textbox(
            L + int(0.07 * W), T + int(top_frac * H),
            int(0.90 * W), int(height_frac * H),
        )
        tf = tb.text_frame
        tf.word_wrap = True
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0
        tf.margin_bottom = 0
        return tf

    # --- Número
    tf = _txt(0.12, 0.26)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = spec["numero"]
    r.font.bold = True
    r.font.size = Pt(spec["size_num"])
    r.font.color.rgb = _rgb(col_dorado)
    r.font.name = "Calibri"

    # --- Título
    tf = _txt(0.42, 0.18)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = spec["titulo"]
    r.font.bold = True
    r.font.size = Pt(spec["size_titulo"])
    r.font.color.rgb = _rgb(col_texto)
    r.font.name = "Calibri"

    # --- Comparativos
    if spec["mostrar_comparativos"]:
        _linea_comparativo(_txt(0.64, 0.15), spec["d25"], "vs 2025 ", spec["size_delta"])
        _linea_comparativo(_txt(0.79, 0.15), spec["d24"], "vs 2024 ", spec["size_delta"])


def _linea_comparativo(tf, delta, etiqueta, size_delta):
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT

    r1 = p.add_run()
    r1.text = f"{delta['icon']}{etiqueta}"
    r1.font.size = Pt(size_delta)
    r1.font.color.rgb = _rgb(col_muted)
    r1.font.name = "Calibri"

    r2 = p.add_run()
    r2.text = delta["label"]
    r2.font.size = Pt(size_delta)
    r2.font.bold = True
    r2.font.color.rgb = _rgb(delta["col"])
    r2.font.name = "Calibri"


# ---------------------------------------------------------------------------
# Helpers de datos (dplyr/tidyr -> pandas)
# ---------------------------------------------------------------------------
def rellenar_anios(df: pd.DataFrame, anios) -> pd.DataFrame:
    """Asegura una fila por año del rango, rellenando numéricos con 0."""
    base = pd.DataFrame({"anio": list(anios)})
    out = base.merge(df, on="anio", how="left")
    num_cols = out.select_dtypes(include=[np.number]).columns
    out[num_cols] = out[num_cols].fillna(0)
    return out


def asegurar_columnas(df: pd.DataFrame, cols) -> pd.DataFrame:
    """Agrega columnas faltantes con valor 0."""
    df = df.copy()
    for c in cols:
        if c not in df.columns:
            df[c] = 0
    return df


def obtener_col(df: pd.DataFrame, nombre, default=0):
    """Devuelve la columna si existe; si no, un vector de `default`."""
    if nombre in df.columns:
        return df[nombre]
    return pd.Series([default] * len(df), index=df.index)


def valor_anio_col(df: pd.DataFrame, columna, anio_objetivo):
    """Primer valor de `columna` para el año dado (NaN si no existe)."""
    if columna not in df.columns:
        return np.nan
    val = df.loc[df["anio"] == anio_objetivo, columna]
    if len(val) == 0:
        return np.nan
    return val.iloc[0]


def tiene_dato_2026(df: pd.DataFrame, columna) -> bool:
    if columna not in df.columns:
        return False
    valor = df.loc[df["anio"] == 2026, columna]
    if len(valor) == 0:
        return False
    return bool(((~valor.isna()) & (valor != 0)).any())


def hay_indicador_2026(df: pd.DataFrame, columna) -> bool:
    val = valor_anio_col(df, columna, 2026)
    return (not pd.isna(val)) and val != 0


# ---------------------------------------------------------------------------
# Value boxes
# ---------------------------------------------------------------------------
def crear_valueboxes_2026(df_3anios, mapa_titulos, sufijo="",
                          incluir_comparativos=True,
                          acento_sin_comparativo="#B0B0B0",
                          size_num=18, size_titulo=11, size_delta=9) -> dict:
    """Crea un dict {clave+sufijo: spec_de_tarjeta} para el año 2026."""
    d26 = df_3anios[df_3anios["anio"] == 2026]

    cards = {}
    for var, titulo in mapa_titulos.items():
        numero = d26[var] if var in d26.columns else pd.Series(dtype=float)

        if len(numero) == 0 or numero.isna().all():
            numero = 0
        else:
            numero = numero.iloc[0]

        if incluir_comparativos:
            col25 = f"var_2026_vs_2025_{var}"
            col24 = f"var_2026_vs_2024_{var}"

            var_vs_2025 = d26[col25].iloc[0] if col25 in d26.columns and len(d26) else np.nan
            var_vs_2024 = d26[col24].iloc[0] if col24 in d26.columns and len(d26) else np.nan

            acento = elige_acento(var_vs_2025, var_vs_2024)
        else:
            var_vs_2025 = None
            var_vs_2024 = None
            acento = acento_sin_comparativo

        cards[f"{var}{sufijo}"] = crear_card_institucional(
            numero=numero, titulo=titulo,
            var_vs_2025=var_vs_2025, var_vs_2024=var_vs_2024,
            acento=acento,
            size_num=size_num, size_titulo=size_titulo, size_delta=size_delta,
        )
    return cards


def definir_layout_valueboxes(datos_consulta_funcion) -> dict:
    """Elige el layout y las métricas según qué indicadores tienen dato 2026."""
    hay_general = tiene_dato_2026(datos_consulta_funcion, "consulta_gral")
    hay_esp = tiene_dato_2026(datos_consulta_funcion, "consulta_esp")
    hay_qx = tiene_dato_2026(datos_consulta_funcion, "qx")
    hay_egresos = tiene_dato_2026(datos_consulta_funcion, "egresos")

    if hay_general and not hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "2_valueboxes", "metricas": ["consulta_gral"]}

    if hay_general and hay_esp and not hay_qx and not hay_egresos:
        return {"layout": "6_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp"]}

    if hay_general and hay_esp and hay_qx and not hay_egresos:
        return {"layout": "8_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx"]}

    if hay_general and hay_esp and hay_qx and hay_egresos:
        return {"layout": "10_valueboxes",
                "metricas": ["total_consultas", "consulta_gral", "consulta_esp", "qx", "egresos"]}

    # Fallback por combinaciones inusuales
    metricas_presentes = []
    if hay_general:
        metricas_presentes.append("consulta_gral")
    if hay_esp:
        metricas_presentes.append("consulta_esp")
    if hay_qx:
        metricas_presentes.append("qx")
    if hay_egresos:
        metricas_presentes.append("egresos")

    if len(metricas_presentes) >= 2:
        metricas_presentes = ["total_consultas"] + metricas_presentes

    n = len(metricas_presentes)
    if n <= 1:
        layout_fallback = "2_valueboxes"
    elif n == 3:
        layout_fallback = "6_valueboxes"
    elif n == 4:
        layout_fallback = "8_valueboxes"
    else:  # n >= 5
        layout_fallback = "10_valueboxes"

    return {"layout": layout_fallback, "metricas": metricas_presentes}


def definir_layout_historico(datos_consulta_funcion):
    hay_consultas = any([
        hay_indicador_2026(datos_consulta_funcion, "total_consultas"),
        hay_indicador_2026(datos_consulta_funcion, "consulta_gral"),
        hay_indicador_2026(datos_consulta_funcion, "consulta_esp"),
    ])
    hay_proc = any([
        hay_indicador_2026(datos_consulta_funcion, "qx"),
        hay_indicador_2026(datos_consulta_funcion, "egresos"),
    ])

    if hay_consultas and hay_proc:
        return "Historico consultas y procedimientos"
    if hay_consultas:
        return "Historico consultas"
    return None


# ---------------------------------------------------------------------------
# Helpers de placeholders / colocación en el .pptx
# ---------------------------------------------------------------------------
def _buscar_layout(prs: Presentation, nombre: str):
    for master in prs.slide_masters:
        for layout in master.slide_layouts:
            if layout.name == nombre:
                return layout
    raise ValueError(f"No se encontró el layout: {nombre!r}")


def _ph_por_nombre(contenedor, nombre):
    for ph in contenedor.placeholders:
        if ph.name == nombre:
            return ph
    return None


def _ph_en_slide_por_nombre(slide, layout, nombre):
    """Busca en `slide` el placeholder correspondiente al `nombre` definido
    en `layout`.

    python-pptx NO conserva el nombre personalizado del placeholder al
    clonarlo del layout hacia la diapositiva (p. ej. "Título 1" del layout
    se vuelve "Title 1" en la diapositiva) — solo conserva el `idx`. Por
    eso se resuelve el idx en el layout y se busca ese idx en la
    diapositiva, en vez de comparar nombres directamente en la diapositiva.
    """
    lp = _ph_por_nombre(layout, nombre)
    if lp is None:
        # Por si acaso el nombre sí se conservó (compatibilidad).
        return _ph_por_nombre(slide, nombre)
    idx = lp.placeholder_format.idx
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def _geom_placeholder(slide, layout, nombre):
    """Devuelve (placeholder_en_slide, (L, T, W, H)) con herencia del layout."""
    ph = _ph_en_slide_por_nombre(slide, layout, nombre)
    lp = _ph_por_nombre(layout, nombre)
    ref = ph if ph is not None else lp
    if ref is None:
        return None, (None, None, None, None)

    def _pick(attr):
        v = getattr(ph, attr) if ph is not None else None
        if v is None and lp is not None:
            v = getattr(lp, attr)
        return v

    return ph, (_pick("left"), _pick("top"), _pick("width"), _pick("height"))


def _quitar_placeholder(ph):
    if ph is not None:
        ph._element.getparent().remove(ph._element)


def _set_texto_placeholder(slide, layout, nombre, texto):
    ph = _ph_en_slide_por_nombre(slide, layout, nombre)
    if ph is not None:
        ph.text = texto


def _dibujar_en_placeholder(slide, layout, nombre, funcion_dibujo, *args, **kwargs):
    """Resuelve el placeholder `nombre`, lo quita, y llama a
    `funcion_dibujo(slide, box, *args, **kwargs)` para dibujar formas
    nativas dentro de su geometria (L, T, W, H)."""
    ph, box = _geom_placeholder(slide, layout, nombre)
    if box[0] is None:
        return
    _quitar_placeholder(ph)
    funcion_dibujo(slide, box, *args, **kwargs)


def _colocar_card(slide, layout, nombre, spec):
    ph, box = _geom_placeholder(slide, layout, nombre)
    if box[0] is None:
        return
    _quitar_placeholder(ph)
    dibujar_card(slide, box, spec)


def imprimir_valueboxes_dinamicos(prs, layout_name, boxes_superior, boxes_inferior,
                                  titulo="Productividad IMSS Bienestar",
                                  fecha=None):
    """Agrega una diapositiva de value boxes y coloca las tarjetas."""
    layout = _buscar_layout(prs, layout_name)
    slide = prs.slides.add_slide(layout)

    _set_texto_placeholder(slide, layout, "Título 1", titulo)
    if fecha is not None:
        _set_texto_placeholder(slide, layout, "fecha", fecha)

    for i, spec in enumerate(boxes_superior, start=1):
        if spec is not None:
            _colocar_card(slide, layout, f"arriba {i}", spec)

    for i, spec in enumerate(boxes_inferior, start=1):
        if spec is not None:
            _colocar_card(slide, layout, f"abajo {i}", spec)

    return slide


# ---------------------------------------------------------------------------
# Formas nativas de bajo nivel (editables en PowerPoint, sin imágenes)
# ---------------------------------------------------------------------------
MESES_EN_ABBR = ["Jan", "Feb", "Mar", "Apr", "May", "Jun",
                  "Jul", "Aug", "Sep", "Oct", "Nov", "Dec"]


def _forma_rect(slide, left, top, width, height, color_hex, transparencia_pct=0):
    """Rectángulo nativo relleno. `transparencia_pct`: 0 (opaco) a 100 (invisible)."""
    shp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    if transparencia_pct:
        srgb = shp.fill._xPr.find(qn("a:solidFill")).find(qn("a:srgbClr"))
        alpha = srgb.makeelement(qn("a:alpha"), {"val": str(int((100 - transparencia_pct) * 1000))})
        srgb.append(alpha)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _forma_texto(slide, left, top, width, height, texto, size=9, color_hex=col_texto,
                 bold=False, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
                 rotation=0, wrap=True, font="Calibri"):
    tb = slide.shapes.add_textbox(int(left), int(top), int(max(width, 1)), int(max(height, 1)))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, linea in enumerate(str(texto).split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = _rgb(color_hex)
        r.font.name = font
    if rotation:
        tb.rotation = rotation
    return tb


def _forma_ovalo(slide, cx, cy, radio, color_hex):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, int(cx - radio), int(cy - radio), int(radio * 2), int(radio * 2))
    shp.fill.solid()
    shp.fill.fore_color.rgb = _rgb(color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _forma_linea_quebrada(slide, puntos, color_hex, width_pt=1.5):
    """Línea poligonal nativa y editable (freeform) que conecta `puntos`
    (lista de (x_emu, y_emu))."""
    if len(puntos) < 2:
        return None
    fb = slide.shapes.build_freeform(int(puntos[0][0]), int(puntos[0][1]), scale=1.0)
    fb.add_line_segments([(int(x), int(y)) for x, y in puntos[1:]], close=False)
    shp = fb.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = _rgb(color_hex)
    shp.line.width = Pt(width_pt)
    shp.shadow.inherit = False
    return shp


def _forma_flecha_vertical(slide, x, y_top, y_bottom, color_hex, width_pt=1.2):
    """Conector vertical con punta de flecha arriba (en `y_top`)."""
    conn = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, int(x), int(y_top), int(x), int(y_bottom))
    conn.line.color.rgb = _rgb(color_hex)
    conn.line.width = Pt(width_pt)
    ln = conn.line._get_or_add_ln()
    head = ln.makeelement(qn("a:headEnd"), {"type": "triangle", "w": "med", "len": "med"})
    ln.append(head)
    return conn


def _escala_bonita(valor_max, n_pasos_objetivo=5):
    """Devuelve (paso, valor_redondeado) con incrementos 'bonitos' (1/2/2.5/5/10 x 10^n),
    igual que el locator automático de matplotlib, para que el eje Y no muestre
    números arbitrarios como 9,420,922."""
    if valor_max <= 0:
        return 1, n_pasos_objetivo
    bruto = valor_max / n_pasos_objetivo
    exponente = math.floor(math.log10(bruto))
    base = 10 ** exponente
    paso = base * 10
    for m in (1, 2, 2.5, 5, 10):
        if bruto <= m * base:
            paso = m * base
            break
    return paso, math.ceil(valor_max / paso) * paso


def _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                   ancho_etiqueta=None, gridlines=False):
    """Dibuja las marcas del eje Y (texto, en `valores_marca`) y opcionalmente
    líneas de rejilla, mapeadas proporcionalmente contra `ymax_eje`."""
    ancho_etiqueta = ancho_etiqueta or Emu(900000)
    for valor in valores_marca:
        frac = (valor / ymax_eje) if ymax_eje else 0
        y = baseline - frac * plot_h
        _forma_texto(
            slide, plot_l - ancho_etiqueta - Pt(4), y - Pt(7), ancho_etiqueta, Pt(14),
            fmt_num(valor), size=8, color_hex=col_muted, align=PP_ALIGN.RIGHT,
            anchor=MSO_ANCHOR.MIDDLE, wrap=False)
        if gridlines and valor > 0:
            _forma_rect(slide, plot_l, y, plot_w, Pt(0.6), "#E5E7EB")


# ---------------------------------------------------------------------------
# Gráficas de barras (histórico 2020-2025 y 2024-2026) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_barras(slide, box, categorias, totales, avances, titulo,
                           colores_total, colores_avance,
                           etiquetas_total=None, etiquetas_avance=None,
                           title_size=15):
    """Barras de 'total' (claro) con 'avance' superpuesto (oscuro), con
    etiquetas de valor — equivalente nativo de grafica_planeacion_*."""
    L, T, W, H = box
    n = len(categorias)
    if n == 0:
        return

    if etiquetas_total is None:
        etiquetas_total = [fmt_num(t) for t in totales]
    if etiquetas_avance is None:
        etiquetas_avance = [fmt_num(a) for a in avances]

    alto_titulo = int(H * 0.13)
    alto_categoria = Pt(16)
    margen_izq = int(W * 0.12)
    espacio_etiqueta_sup = int(H * 0.16)

    _forma_texto(slide, L, T, W, alto_titulo, titulo, size=title_size,
                color_hex=col_muted, bold=True, align=PP_ALIGN.CENTER)

    plot_l = L + margen_izq
    plot_t = T + alto_titulo + espacio_etiqueta_sup
    plot_w = W - margen_izq
    plot_h = H - alto_titulo - espacio_etiqueta_sup - alto_categoria
    baseline = plot_t + plot_h

    ymax_datos = max(max(totales, default=0), max(avances, default=0), 1)
    paso, ymax_redondeado = _escala_bonita(ymax_datos, 5)
    ymax_eje = max(ymax_redondeado, ymax_datos) * 1.16

    valores_marca = [i * paso for i in range(int(ymax_eje // paso) + 1)]
    _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca,
                  ancho_etiqueta=margen_izq - Pt(4))

    slot_w = plot_w / n
    bar_w = slot_w * 0.60

    for i, cat in enumerate(categorias):
        slot_l = plot_l + i * slot_w
        bar_l = slot_l + (slot_w - bar_w) / 2

        h_total = (totales[i] / ymax_eje) * plot_h if ymax_eje else 0
        h_avance = (avances[i] / ymax_eje) * plot_h if ymax_eje else 0

        if h_total > 0:
            _forma_rect(slide, bar_l, baseline - h_total, bar_w, h_total, colores_total[i])
        if h_avance > 0:
            _forma_rect(slide, bar_l, baseline - h_avance, bar_w, h_avance, colores_avance[i])

        n_lineas_tot = etiquetas_total[i].count("\n") + 1
        alto_et = Pt(13) * n_lineas_tot
        _forma_texto(slide, slot_l, baseline - h_total - alto_et - Pt(3), slot_w, alto_et,
                    etiquetas_total[i], size=9.5, color_hex="#000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

        if h_avance > Pt(24):
            n_lineas_av = etiquetas_avance[i].count("\n") + 1
            alto_ea = Pt(12) * n_lineas_av
            _forma_texto(slide, slot_l, baseline - h_avance + Pt(3), slot_w, alto_ea,
                        etiquetas_avance[i], size=8.5, color_hex="#FFFFFF", bold=True,
                        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

        _forma_texto(slide, slot_l, baseline + Pt(3), slot_w, alto_categoria,
                    str(cat), size=11, color_hex=col_muted, bold=True, align=PP_ALIGN.CENTER)


def dibujar_grafica_planeacion_historica(slide, box, df, col_total, col_avance, titulo,
                                         beige="#D9D2BE", verde="#2F6F63"):
    """Barras 2020-2025: total (beige) con avance (verde) superpuesto."""
    anios = list(range(2020, 2026))
    d = df.copy()
    d["anio_num"] = pd.to_numeric(d["anio"], errors="coerce").astype("Int64")
    d = d[d["anio_num"].isin(anios)]
    d = d.set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=[beige] * len(anios), colores_avance=[verde] * len(anios),
    )


def dibujar_grafica_planeacion_2024_2026(slide, box, df, col_total, col_avance, titulo,
                                         beige="#D9D2BE", verde="#2F6F63",
                                         beige_2026="#A99F86", verde_2026="#1E5B4F"):
    """Barras 2024-2026 con etiquetas especiales de 'Meta 2026' / 'Avance'."""
    anios = [2024, 2025, 2026]
    d = df.copy()
    d["anio_num"] = pd.to_numeric(d["anio"], errors="coerce").astype("Int64")
    d = d[d["anio_num"].isin(anios)]
    d = d.set_index("anio_num").reindex(anios)
    totales = d[col_total].fillna(0).astype(float).tolist()
    avances = d[col_avance].fillna(0).astype(float).tolist()

    etiquetas_total, etiquetas_avance = [], []
    for a, tot, av in zip(anios, totales, avances):
        if a == 2026:
            etiquetas_total.append(f"Meta 2026\n{fmt_num(tot)}")
            pct = (av / tot) if tot > 0 else np.nan
            pct_txt = "s/d" if pd.isna(pct) else f"{int(round(pct * 100))}%"
            etiquetas_avance.append(f"Avance\n{fmt_num(av)}\n({pct_txt})")
        else:
            etiquetas_total.append(fmt_num(tot))
            etiquetas_avance.append(fmt_num(av))

    colores_total = [beige_2026 if a == 2026 else beige for a in anios]
    colores_avance = [verde_2026 if a == 2026 else verde for a in anios]

    dibujar_grafica_barras(
        slide, box, categorias=[str(a) for a in anios],
        totales=totales, avances=avances, titulo=titulo,
        colores_total=colores_total, colores_avance=colores_avance,
        etiquetas_total=etiquetas_total, etiquetas_avance=etiquetas_avance,
    )


# ---------------------------------------------------------------------------
# Gráfica de serie temporal (consultas/procedimientos por mes) — formas nativas
# ---------------------------------------------------------------------------
def dibujar_grafica_consultas_periodos(slide, box, df, fecha_inicio="2022-08-01", fecha_fin=None,
                                       color_linea="#6B6B6B", verde_punto="#1F5B50",
                                       fill_2223="#EFEFEF", fill_2024="#E9DDCC",
                                       fill_2025="#F4F0EA", fill_2026="#E9DDCC",
                                       fill_valuebox="#B99C6D"):
    """Serie temporal mensual con bandas por periodo y anotaciones, dibujada
    con formas nativas de PowerPoint (línea freeform editable, rectángulos,
    óvalos y textos) en vez de una imagen."""
    L, T, W, H = box
    hoy_mes = _floor_month(date.today())
    if fecha_fin is None:
        fecha_fin = hoy_mes
    else:
        fecha_fin = pd.to_datetime(fecha_fin).date()
    fecha_inicio = pd.to_datetime(fecha_inicio).date()

    d = df.copy()
    d["fecha"] = pd.to_datetime(d["fecha"]).dt.tz_localize(None)
    d = d[(d["fecha"].dt.date >= fecha_inicio) &
          (d["fecha"].dt.date <= fecha_fin) &
          (d["fecha"].dt.date < hoy_mes)].sort_values("fecha")

    if len(d) == 0:
        _forma_texto(slide, L, T, W, H, "Sin datos suficientes para este período",
                    size=12, color_hex=col_muted)
        return

    ymax = float(d["consultas_totales"].max())
    ymin = float(d["consultas_totales"].min())
    ymax_eje = max(ymax * 1.48, 1)
    _paso_eje, _ymax_marcas = _escala_bonita(ymax, 5)
    valores_marca_y = [i * _paso_eje for i in range(int(ymax_eje // _paso_eje) + 1)]

    fecha_fin_banda = _ceiling_month(fecha_fin)
    bandas = [
        (fecha_inicio, date(2024, 1, 1), fill_2223, "2022–2023\nAños de transición"),
        (date(2024, 1, 1), date(2025, 1, 1), fill_2024, "2024\nPrimer año de operación"),
        (date(2025, 1, 1), date(2026, 1, 1), fill_2025, "2025\nSegundo año de operación"),
        (date(2026, 1, 1), fecha_fin_banda, fill_2026, "2026\nTercer año de operación"),
    ]

    dias_pad = max((fecha_fin_banda - fecha_inicio).days * 0.035, 10)
    x0_ord = fecha_inicio.toordinal() - dias_pad
    x1_ord = fecha_fin_banda.toordinal() + dias_pad

    margen_izq = int(W * 0.075)
    margen_der = int(W * 0.015)
    margen_sup = int(H * 0.03)
    margen_inf = int(H * 0.17)

    plot_l = L + margen_izq
    plot_t = T + margen_sup
    plot_w = W - margen_izq - margen_der
    plot_h = H - margen_sup - margen_inf
    baseline = plot_t + plot_h

    def _ord(valor_fecha):
        if hasattr(valor_fecha, "toordinal"):
            return valor_fecha.toordinal()
        return pd.Timestamp(valor_fecha).toordinal()

    def xmap(valor_fecha):
        frac = (_ord(valor_fecha) - x0_ord) / (x1_ord - x0_ord)
        return plot_l + frac * plot_w

    def ymap(valor):
        frac = (valor / ymax_eje) if ymax_eje else 0
        return baseline - frac * plot_h

    # Fondo blanco del área de la gráfica
    _forma_rect(slide, L, T, W, H, "#FFFFFF")

    # Bandas de periodo
    for xmin_b, xmax_b, fill, _lab in bandas:
        x_l = xmap(xmin_b)
        x_r = xmap(xmax_b)
        _forma_rect(slide, x_l, plot_t, x_r - x_l, plot_h, fill)

    # Rejilla horizontal + etiquetas del eje Y
    _dibujar_eje_y(slide, plot_l, plot_w, baseline, plot_h, ymax_eje, valores_marca_y,
                  ancho_etiqueta=margen_izq - Pt(4), gridlines=True)

    # Zona posible subregistro (últimos 3 meses)
    ult3 = d.tail(3)
    if len(ult3):
        xmin_sr = (ult3["fecha"].min() - pd.Timedelta(days=15)).date()
        xmax_sr = (ult3["fecha"].max() + pd.Timedelta(days=15)).date()
        x_l = xmap(xmin_sr)
        x_r = xmap(xmax_sr)
        _forma_rect(slide, x_l, plot_t, x_r - x_l, plot_h, "#B22222", transparencia_pct=82)
        _forma_texto(
            slide, xmap(d["fecha"].max().date()) - Pt(70), ymap(ymax * 1.08), Pt(140), Pt(26),
            "Posible subregistro\ntemporal", size=8.5, color_hex="#7A1E3A", bold=True,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Serie: línea nativa (freeform, editable) + marcador en cada punto
    puntos_linea = [(xmap(f.date()), ymap(v)) for f, v in zip(d["fecha"], d["consultas_totales"])]
    _forma_linea_quebrada(slide, puntos_linea, color_linea, width_pt=1.3)
    for x, y in puntos_linea:
        _forma_ovalo(slide, x, y, Pt(1.6), color_linea)

    # Puntos destacados (mismo mes que el corte, años anteriores a 2026)
    mes_destacado = fecha_fin.month
    puntos_destacados = d[(d["fecha"].dt.month == mes_destacado) &
                          (d["fecha"].dt.year < 2026)].drop_duplicates("fecha")
    for _, row in puntos_destacados.iterrows():
        x = xmap(row["fecha"].date())
        y = ymap(row["consultas_totales"])
        _forma_ovalo(slide, x, y, Pt(4.5), verde_punto)
        etiqueta = f"{fmt_num(row['consultas_totales'])}\n{_mes_abbr_title(row['fecha'])}"
        _forma_texto(slide, x - Pt(45), y - Pt(38), Pt(90), Pt(30), etiqueta,
                    size=8, color_hex="#000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

    # Texto superior de bandas
    for xmin_b, xmax_b, _fill, lab in bandas:
        centro = xmap(xmin_b) + (xmap(xmax_b) - xmap(xmin_b)) / 2
        _forma_texto(slide, centro - Pt(70), ymap(ymax * 1.32), Pt(140), Pt(26), lab,
                    size=8, color_hex="#000000", bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

    # Flecha "Decreto de creación"
    x_decreto = xmap(date(2022, 8, 15))
    _forma_flecha_vertical(slide, x_decreto, ymap(ymax * 1.02), ymap(ymin * 0.95), verde_punto)
    _forma_texto(slide, x_decreto + Pt(6), ymap(ymax * 1.05) - Pt(4), Pt(120), Pt(26),
                "Decreto de creación\ndel IMSS Bienestar", size=7.5, color_hex="#000000",
                bold=True, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)

    # Recuadro (valuebox) del último dato
    fecha_ultimo = d["fecha"].max()
    valor_ultimo_ser = d.loc[d["fecha"] == fecha_ultimo, "consultas_totales"]
    valor_ultimo = valor_ultimo_ser.iloc[0] if len(valor_ultimo_ser) else 0
    if pd.isna(valor_ultimo):
        valor_ultimo = 0
    x_vb = xmap(fecha_ultimo.date())
    y_vb = ymap(valor_ultimo)
    vb_w, vb_h = Pt(72), Pt(30)
    vb = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, int(x_vb - vb_w / 2), int(y_vb - vb_h - Pt(10)),
        int(vb_w), int(vb_h))
    vb.fill.solid()
    vb.fill.fore_color.rgb = _rgb(fill_valuebox)
    vb.line.fill.background()
    vb.shadow.inherit = False
    tf = vb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for i, linea in enumerate([fmt_num(valor_ultimo), _mes_title(fecha_ultimo)]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = linea
        r.font.size = Pt(9.5)
        r.font.bold = True
        r.font.color.rgb = _rgb("#FFFFFF")
        r.font.name = "Calibri"

    # Etiquetas del eje X (cada 2 meses)
    eje_inicio = pd.Timestamp(fecha_inicio).replace(day=1) - pd.DateOffset(months=2)
    eje_fin = pd.Timestamp(fecha_fin_banda) + pd.DateOffset(months=2)
    for tick in pd.date_range(eje_inicio, eje_fin, freq="2MS"):
        x = xmap(tick.date())
        if x < plot_l - Pt(5) or x > plot_l + plot_w + Pt(5):
            continue
        etiqueta = f"{MESES_EN_ABBR[tick.month - 1]}-{str(tick.year)[2:]}"
        _forma_texto(slide, x - Pt(20), baseline + Pt(4), Pt(40), Pt(22), etiqueta,
                    size=7.5, color_hex=col_muted, bold=False, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.TOP, rotation=45, wrap=False)


# ---------------------------------------------------------------------------
# Tablas (flextable -> tabla nativa de python-pptx)
# ---------------------------------------------------------------------------
def armar_tabla_dinamica(df, indicadores, etiquetas, mes_nombre) -> pd.DataFrame:
    """Construye la tabla comparativa 2025 vs 2026 con crecimiento anual."""
    valor_2025 = [valor_anio_col(df, ind, 2025) for ind in indicadores]
    valor_2026 = [valor_anio_col(df, ind, 2026) for ind in indicadores]

    filas = []
    for etq, v25, v26 in zip(etiquetas, valor_2025, valor_2026):
        if pd.isna(v25) or v25 == 0:
            crecimiento = np.nan
        else:
            crecimiento = round((v26 - v25) / v25 * 100, 0)

        if pd.isna(crecimiento):
            crec_txt = "s/d"
        elif crecimiento > 0:
            crec_txt = f"+{int(crecimiento)} %"
        else:
            crec_txt = f"{int(crecimiento)} %"

        v25_txt = "s/d" if pd.isna(v25) else fmt_num(v25)
        v26_txt = "s/d" if pd.isna(v26) else fmt_num(v26)

        filas.append({
            "indicador": etq,
            f"{mes_nombre} 2025": v25_txt,
            f"{mes_nombre} 2026": v26_txt,
            "Crecimiento anual": crec_txt,
        })

    return pd.DataFrame(filas, columns=[
        "indicador", f"{mes_nombre} 2025", f"{mes_nombre} 2026", "Crecimiento anual"])


def _set_cell_border(cell, color="6B7280", width_pt=1):
    """Agrega borde a las 4 aristas de una celda (python-pptx no lo expone)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        tag = qn(edge)
        for existing in tcPr.findall(tag):
            tcPr.remove(existing)
        ln = tcPr.makeelement(tag, {
            "w": str(int(Pt(width_pt))), "cap": "flat",
            "cmpd": "sng", "algn": "ctr",
        })
        fill = ln.makeelement(qn("a:solidFill"), {})
        clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
        fill.append(clr)
        ln.append(fill)
        tcPr.append(ln)


def ft_planeacion(prs_slide, layout, nombre_ph, df,
                  w1=4.05, w2=1.20, w3=1.20, w4=1.35,
                  header_negro="#3B3B3B", verde="#2F6F63", menta="#D9F2EE",
                  size_header=None, size_body=None, h_fila=None):
    """Dibuja una tabla estilizada en el placeholder indicado."""
    n_filas = len(df)

    if size_header is None or size_body is None or h_fila is None:
        if n_filas == 1:
            size_header, size_body, h_fila = 13, 12, 0.45
        elif n_filas == 2:
            size_header, size_body, h_fila = 11, 10, 0.35
        else:
            size_header, size_body, h_fila = 10, 9, 0.28

    ph, (L, T, W, H) = _geom_placeholder(prs_slide, layout, nombre_ph)
    if L is None:
        return
    _quitar_placeholder(ph)

    encabezados = ["Indicador"] + list(df.columns[1:])
    n_cols = len(encabezados)
    anchos = [w1, w2, w3, w4][:n_cols]
    ancho_total = Inches(sum(anchos))
    alto_total = Inches(h_fila * (n_filas + 1))

    gf = prs_slide.shapes.add_table(n_filas + 1, n_cols, L, T, ancho_total, alto_total)
    tabla = gf.table
    tabla.first_row = False
    tabla.horz_banding = False

    for j, wi in enumerate(anchos):
        tabla.columns[j].width = Inches(wi)
    for i in range(n_filas + 1):
        tabla.rows[i].height = Inches(h_fila)

    # Encabezado
    for j, texto in enumerate(encabezados):
        cell = tabla.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = _rgb(header_negro)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
        r = p.add_run()
        r.text = str(texto)
        r.font.bold = True
        r.font.size = Pt(size_header)
        r.font.color.rgb = _rgb("#FFFFFF")
        _set_cell_border(cell)

    # Cuerpo
    for i in range(n_filas):
        for j, colname in enumerate(df.columns):
            cell = tabla.cell(i + 1, j)
            if j == 2:  # tercera columna (año 2026) -> fondo menta
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(menta)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb("#FFFFFF")
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if j == 0 else PP_ALIGN.CENTER
            r = p.add_run()
            r.text = str(df.iloc[i, j])
            r.font.size = Pt(size_body)
            r.font.color.rgb = _rgb(col_texto)
            _set_cell_border(cell)

    return tabla


# ---------------------------------------------------------------------------
# Utilidades de fecha
# ---------------------------------------------------------------------------
def _floor_month(d) -> date:
    d = pd.to_datetime(d).date() if not isinstance(d, date) else d
    return date(d.year, d.month, 1)


def _ceiling_month(d) -> date:
    d = pd.to_datetime(d).date() if not isinstance(d, date) else d
    if d.month == 12:
        return date(d.year + 1, 1, 1)
    return date(d.year, d.month + 1, 1)


def _mes_title(ts) -> str:
    ts = pd.to_datetime(ts)
    return f"{MESES_ES[ts.month - 1].capitalize()} {ts.year}"


def _mes_abbr_title(ts) -> str:
    ts = pd.to_datetime(ts)
    abbr = MESES_ES[ts.month - 1][:3].capitalize()
    return f"{abbr}-{ts.year}"


def _calcular_fecha_corte(hoy: date) -> date:
    """Reproduce la lógica de corte semanal (miércoles) del código R."""
    # POSIXlt$wday: domingo=0 ... sábado=6
    wday0 = (hoy.weekday() + 1) % 7
    if wday0 == 3:  # miércoles (lubridate::wday == 4)
        return hoy - timedelta(days=7)
    return hoy - timedelta(days=((wday0 + 4) % 7))


# ---------------------------------------------------------------------------
# FUNCIÓN PRINCIPAL
# ---------------------------------------------------------------------------
_RUTA_MASTER_DEFAULT = (
    Path(__file__).resolve().parent.parent / "inst" / "app" / "data" / "master_presentacion.pptx"
)


def crear_reporte_productividad(codigo_clues, clues_info, metas, historicos,
                                procedimientos_personas=None,
                                ruta_master=None,
                                hoy=None):
    """Genera la presentación de productividad y devuelve el objeto Presentation."""
    if hoy is None:
        hoy = date.today()
    if ruta_master is None:
        ruta_master = str(_RUTA_MASTER_DEFAULT)

    # Variables de corte ----------------------------------------------------
    fecha_corte = _calcular_fecha_corte(hoy)
    fecha_fin_graf = _floor_month(fecha_corte)

    fecha_portada = f"{fecha_corte.day:02d} de {MESES_ES[fecha_corte.month - 1]} de {fecha_corte.year}"
    mes_nombre = MESES_ES[fecha_corte.month - 1].capitalize()

    # Filtrado por plan / nacional -----------------------------------------
    if codigo_clues == "NACIONAL":
        clues_info_filtrado = pd.DataFrame(
            [{"clues": "NACIONAL", "nombre": "NACIONAL", "entidad": "NACIONAL"}])
        metas_filtrado = metas.copy()
    else:
        mask = (clues_info["clues"] == codigo_clues)
        if "id" in clues_info.columns:
            mask = mask | (clues_info["id"] == codigo_clues)
        clues_info_filtrado = clues_info[mask].copy()
        metas_filtrado = metas[metas["clues"] == codigo_clues].copy()

    if len(metas_filtrado) == 0:
        raise ValueError(f"No hay metas para: {codigo_clues}")

    if len(clues_info_filtrado) == 0:
        clues_info_filtrado = pd.DataFrame(
            [{"clues": codigo_clues, "nombre": codigo_clues, "entidad": np.nan}])

    meta_total_consultas = (
        pd.to_numeric(metas_filtrado.get("meta_general_anual"), errors="coerce").fillna(0).sum()
        + pd.to_numeric(metas_filtrado.get("meta_especialidad_anual"), errors="coerce").fillna(0).sum()
    )
    meta_qx = pd.to_numeric(metas_filtrado.get("meta_cirugia_anual"), errors="coerce").fillna(0).sum()

    # Derivar procedimientos_personas si no viene ---------------------------
    if procedimientos_personas is None:
        h = historicos.copy()
        h["anio"] = pd.to_datetime(h["fecha"]).dt.year
        agg = h.groupby("anio", as_index=False).agg(
            total_consultas=("consulta_total", "sum"),
            consulta_gral=("consulta_general", "sum"),
            consulta_esp=("consulta_especialidad", "sum"),
            qx=("procedimientos_qx", "sum"),
            egresos=("egresos", "sum"),
        )
        largo = agg.melt(id_vars="anio",
                         value_vars=["total_consultas", "consulta_gral",
                                     "consulta_esp", "qx", "egresos"],
                         var_name="tipo_procedimiento", value_name="procedimientos")
        largo["fecha"] = largo["anio"]
        largo["personas"] = largo["procedimientos"]
        procedimientos_personas = largo[["fecha", "tipo_procedimiento",
                                         "procedimientos", "personas"]]

    procedimientos_personas = procedimientos_personas.copy()
    recodificar = {
        "consulta total": "total_consultas",
        "general": "consulta_gral",
        "especialidad": "consulta_esp",
        "qx": "qx",
        "egresos": "egresos",
    }
    procedimientos_personas["tipo_procedimiento"] = (
        procedimientos_personas["tipo_procedimiento"].map(lambda v: recodificar.get(v, v)))

    # Presentación ----------------------------------------------------------
    prs = Presentation(ruta_master)

    # Portada ---------------------------------------------------------------
    layout_portada = _buscar_layout(prs, "Portada 3")
    slide = prs.slides.add_slide(layout_portada)
    titulo_portada = f"Reporte de {str(clues_info_filtrado['clues'].iloc[0]).title()}"
    _set_texto_placeholder(slide, layout_portada, "Título 1", titulo_portada)
    _set_texto_placeholder(slide, layout_portada, "Marcador de contenido 2", fecha_portada)

    # Datos base ------------------------------------------------------------
    cols_metricas = ["consulta_gral", "consulta_esp", "qx", "total_consultas", "egresos"]

    h = historicos.copy()
    h["anio"] = pd.to_datetime(h["fecha"]).dt.year
    # obtener_col: columnas faltantes se tratan como 0
    h["_cgral"] = obtener_col(h, "consulta_general")
    h["_cesp"] = obtener_col(h, "consulta_especialidad")
    h["_qx"] = obtener_col(h, "procedimientos_qx")
    h["_tot"] = obtener_col(h, "consulta_total")
    h["_egr"] = obtener_col(h, "egresos")
    datos_anual = h.groupby("anio", as_index=False).agg(
        consulta_gral_anual=("_cgral", "sum"),
        consulta_esp_anual=("_cesp", "sum"),
        qx_anual=("_qx", "sum"),
        total_consultas_anual=("_tot", "sum"),
        egresos_anual=("_egr", "sum"),
    )
    datos_anual = rellenar_anios(datos_anual, range(2020, 2027))

    # datos_consulta_funcion (procedimientos)
    dcf = procedimientos_personas.copy()
    dcf["anio"] = pd.to_numeric(dcf["fecha"], errors="coerce")
    dcf = (dcf[["anio", "tipo_procedimiento", "procedimientos"]]
           .pivot_table(index="anio", columns="tipo_procedimiento",
                        values="procedimientos", aggfunc="sum")
           .reset_index())
    dcf.columns.name = None
    dcf = rellenar_anios(dcf, range(2024, 2027))
    dcf = asegurar_columnas(dcf, cols_metricas)
    datos_consulta_funcion = dcf.merge(datos_anual, on="anio", how="left").sort_values("anio")
    num_cols = datos_consulta_funcion.select_dtypes(include=[np.number]).columns
    datos_consulta_funcion[num_cols] = datos_consulta_funcion[num_cols].fillna(0)
    datos_consulta_funcion["total_consultas_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026,
        meta_total_consultas, datos_consulta_funcion["total_consultas_anual"])
    datos_consulta_funcion["qx_meta"] = np.where(
        datos_consulta_funcion["anio"] == 2026,
        meta_qx, datos_consulta_funcion["qx_anual"])

    # datos_curps (personas)
    dcu = procedimientos_personas.copy()
    dcu["anio"] = pd.to_numeric(dcu["fecha"], errors="coerce")
    dcu = (dcu[["anio", "tipo_procedimiento", "personas"]]
           .pivot_table(index="anio", columns="tipo_procedimiento",
                        values="personas", aggfunc="sum")
           .reset_index())
    dcu.columns.name = None
    dcu = rellenar_anios(dcu, range(2024, 2027))
    dcu = asegurar_columnas(dcu, cols_metricas)
    datos_curps = dcu
    num_cols = datos_curps.select_dtypes(include=[np.number]).columns
    datos_curps[num_cols] = datos_curps[num_cols].fillna(0)

    # Comparativos ----------------------------------------------------------
    for col in cols_metricas:
        for ref in (2024, 2025):
            nombre_var = f"var_2026_vs_{ref}_{col}"

            v26 = valor_anio_col(datos_consulta_funcion, col, 2026)
            vref = valor_anio_col(datos_consulta_funcion, col, ref)
            cond = (not pd.isna(v26)) and (not pd.isna(vref)) and vref != 0
            valor = round(100 * (1 - (vref / v26)), 0) if cond and v26 not in (0, np.nan) else 0
            datos_consulta_funcion[nombre_var] = np.where(
                datos_consulta_funcion["anio"] == 2026, valor, 0)

            v26c = valor_anio_col(datos_curps, col, 2026)
            vrefc = valor_anio_col(datos_curps, col, ref)
            condc = (not pd.isna(v26c)) and (not pd.isna(vrefc)) and vrefc != 0
            valorc = round(100 * (1 - (vrefc / v26c)), 0) if condc and v26c not in (0, np.nan) else 0
            datos_curps[nombre_var] = np.where(datos_curps["anio"] == 2026, valorc, 0)

    # Value boxes -----------------------------------------------------------
    mapa_titulos_consultas = {
        "total_consultas": "Consultas totales",
        "consulta_gral": "Consulta general",
        "consulta_esp": "Especialidad",
        "qx": "Procedimientos quirúrgicos",
        "egresos": "Egresos",
    }
    mapa_titulos_curp = {
        "total_consultas": "Consultas totales",
        "consulta_gral": "Consulta general",
        "consulta_esp": "Especialidad",
        "qx": "Intervenidas",
        "egresos": "Egresadas",
    }

    vbox_consultas = crear_valueboxes_2026(
        datos_consulta_funcion, mapa_titulos_consultas, incluir_comparativos=True)
    vbox_curps = crear_valueboxes_2026(
        datos_curps, mapa_titulos_curp, sufijo="_p", incluir_comparativos=True)

    config_vb = definir_layout_valueboxes(datos_consulta_funcion)
    layout_vb = config_vb["layout"]
    metricas_vb = config_vb["metricas"]

    lista_boxes_superior = {
        "total_consultas": vbox_consultas.get("total_consultas"),
        "consulta_gral": vbox_consultas.get("consulta_gral"),
        "consulta_esp": vbox_consultas.get("consulta_esp"),
        "qx": vbox_consultas.get("qx"),
        "egresos": vbox_consultas.get("egresos"),
    }
    lista_boxes_inferior = {
        "total_consultas": vbox_curps.get("total_consultas_p"),
        "consulta_gral": vbox_curps.get("consulta_gral_p"),
        "consulta_esp": vbox_curps.get("consulta_esp_p"),
        "qx": vbox_curps.get("qx_p"),
        "egresos": vbox_curps.get("egresos_p"),
    }

    imprimir_valueboxes_dinamicos(
        prs=prs,
        layout_name=layout_vb,
        boxes_superior=[lista_boxes_superior[m] for m in metricas_vb],
        boxes_inferior=[lista_boxes_inferior[m] for m in metricas_vb],
        titulo="Productividad IMSS Bienestar",
        fecha=f"Del 01 de enero al {fecha_portada}",
    )

    # Diapo 3: histórico 2020-2025 -----------------------------------------
    fecha_corte_15 = date(fecha_corte.year, fecha_corte.month, 15)
    corte_md = f"{fecha_corte_15.month:02d}-{fecha_corte_15.day:02d}"

    hh = historicos.copy()
    hh["fecha_dt"] = pd.to_datetime(hh["fecha"]).dt.tz_localize(None)
    hh["consulta_general_tmp"] = obtener_col(hh, "consulta_general") + obtener_col(hh, "consulta_gral")
    hh["consulta_esp_tmp"] = obtener_col(hh, "consulta_especialidad") + obtener_col(hh, "consulta_esp")
    hh["procedimientos_qx_tmp"] = obtener_col(hh, "procedimientos_qx")
    hh["egresos_tmp"] = obtener_col(hh, "egresos")
    hh["md"] = hh["fecha_dt"].dt.strftime("%m-%d")
    hh = hh[(hh["fecha_dt"].dt.year >= 2020) &
            (hh["fecha_dt"].dt.year <= 2025) &
            (hh["md"] <= corte_md)].copy()
    hh["anio"] = hh["fecha_dt"].dt.year.astype(str)

    datos_historicos_2020_2025 = hh.groupby("anio", as_index=False).agg(
        consulta_gral=("consulta_general_tmp", "sum"),
        consulta_esp=("consulta_esp_tmp", "sum"),
        qx=("procedimientos_qx_tmp", "sum"),
        egresos=("egresos_tmp", "sum"),
    )

    def _total_consultas(g, e):
        if g > 0 and e > 0:
            return g + e
        if g > 0:
            return g
        if e > 0:
            return e
        return 0

    datos_historicos_2020_2025["total_consultas"] = [
        _total_consultas(g, e) for g, e in
        zip(datos_historicos_2020_2025["consulta_gral"],
            datos_historicos_2020_2025["consulta_esp"])]

    da_str = datos_anual.copy()
    da_str["anio"] = da_str["anio"].astype(str)
    datos_historicos_2020_2025 = datos_historicos_2020_2025.merge(
        da_str, on="anio", how="left")
    for c in ["consulta_gral_anual", "consulta_esp_anual", "qx_anual"]:
        datos_historicos_2020_2025[c] = datos_historicos_2020_2025.get(c, 0)
        datos_historicos_2020_2025[c] = datos_historicos_2020_2025[c].fillna(0)
    datos_historicos_2020_2025["total_consultas_anual"] = [
        _total_consultas(g, e) for g, e in
        zip(datos_historicos_2020_2025["consulta_gral_anual"],
            datos_historicos_2020_2025["consulta_esp_anual"])]
    datos_historicos_2020_2025 = datos_historicos_2020_2025.sort_values("anio")

    print("REVISION DATOS HISTORICOS")
    print(datos_historicos_2020_2025)

    hay_consultas_2020 = bool((datos_historicos_2020_2025["total_consultas"] > 0).any())
    hay_qx_2020 = bool(((datos_historicos_2020_2025["qx"] > 0) |
                        (datos_historicos_2020_2025["egresos"] > 0)).any())

    if hay_consultas_2020:
        if hay_qx_2020:
            lay = _buscar_layout(prs, "1_Historico consultas y procedimientos")
            s = prs.slides.add_slide(lay)
            _set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            _dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "total_consultas_anual",
                                    "total_consultas", "Consultas totales")
            _dibujar_en_placeholder(s, lay, "Grafica 2", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "qx_anual", "qx",
                                    "Procedimientos quirúrgicos")
            _set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            lay = _buscar_layout(prs, "1_Historico consultas")
            s = prs.slides.add_slide(lay)
            _set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            _dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_historica,
                                    datos_historicos_2020_2025, "total_consultas_anual",
                                    "total_consultas", "Consultas totales")
            _set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 4: 2024-2026 con tablas ----------------------------------------
    for col in ("consulta_gral", "consulta_esp", "qx", "egresos"):
        if col not in datos_consulta_funcion.columns:
            datos_consulta_funcion[col] = 0

    d2426 = datos_consulta_funcion.copy()
    d2426["anio_num"] = d2426["anio"].astype(int)
    for c in ("consulta_gral", "consulta_esp", "qx", "egresos",
              "consulta_gral_anual", "consulta_esp_anual", "qx_anual"):
        if c in d2426.columns:
            d2426[c] = d2426[c].fillna(0)

    d2426["total_consultas"] = [
        _total_consultas(g, e) for g, e in
        zip(d2426["consulta_gral"], d2426["consulta_esp"])]

    def _total_meta(row):
        if row["anio_num"] == 2026 and not pd.isna(row.get("total_consultas_meta")):
            return row["total_consultas_meta"]
        if row["anio_num"] == 2026:
            return row["total_consultas"]
        g, e = row.get("consulta_gral_anual", 0), row.get("consulta_esp_anual", 0)
        if g > 0 and e > 0:
            return g + e
        if g > 0:
            return g
        if e > 0:
            return e
        if row.get("total_consultas_anual", 0) > 0:
            return row["total_consultas_anual"]
        return row["total_consultas"]

    def _qx_meta(row):
        if row["anio_num"] == 2026 and not pd.isna(row.get("qx_meta")):
            return row["qx_meta"]
        if row["anio_num"] == 2026:
            return row["qx"]
        if row.get("qx_anual", 0) > 0:
            return row["qx_anual"]
        return row["qx"]

    d2426["total_consultas_meta"] = d2426.apply(_total_meta, axis=1)
    d2426["qx_meta"] = d2426.apply(_qx_meta, axis=1)
    d2426 = d2426[d2426["anio_num"].isin([2024, 2025, 2026])]

    print("REVISION DATOS 2024_2026")
    print(d2426)

    hay_consultas_2426 = bool((d2426["total_consultas"] > 0).any())
    hay_qx_2426 = bool(((d2426["qx"] > 0) | (d2426["egresos"] > 0)).any())

    if hay_consultas_2426:
        indicadores_consulta, etiquetas_consulta = [], []
        if hay_indicador_2026(d2426, "consulta_gral"):
            indicadores_consulta.append("consulta_gral")
            etiquetas_consulta.append("Consultas generales")
        if hay_indicador_2026(d2426, "consulta_esp"):
            indicadores_consulta.append("consulta_esp")
            etiquetas_consulta.append("Consultas de especialidad*")

        tabla_consultas = armar_tabla_dinamica(
            d2426, indicadores_consulta, etiquetas_consulta, "Acumulado")

        if hay_qx_2426:
            indicadores_proc, etiquetas_proc = [], []
            if hay_indicador_2026(d2426, "qx"):
                indicadores_proc.append("qx")
                etiquetas_proc.append("Procedimientos quirúrgicos")
            if hay_indicador_2026(d2426, "egresos"):
                indicadores_proc.append("egresos")
                etiquetas_proc.append("Egresos")

            tabla_proc = armar_tabla_dinamica(
                d2426, indicadores_proc, etiquetas_proc, "Acumulado")

            lay = _buscar_layout(prs, "Historico consultas y procedimientos")
            s = prs.slides.add_slide(lay)
            _set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            _dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "total_consultas_meta", "total_consultas",
                                    "Consultas totales")
            _dibujar_en_placeholder(s, lay, "Grafica 2", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "qx_meta", "qx", "Procedimientos quirúrgicos")
            ft_planeacion(s, lay, "tabla_1", tabla_consultas,
                          w1=2.70, w2=0.90, w3=0.90, w4=0.80,
                          size_header=8, size_body=7.5, h_fila=0.28)
            ft_planeacion(s, lay, "tabla_2", tabla_proc,
                          w1=2.70, w2=0.90, w3=0.90, w4=0.80,
                          size_header=8, size_body=7.5, h_fila=0.28)
            _set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")
        else:
            lay = _buscar_layout(prs, "Historico consultas")
            s = prs.slides.add_slide(lay)
            _set_texto_placeholder(s, lay, "Título 1", "Productividad IMSS Bienestar")
            _dibujar_en_placeholder(s, lay, "Grafica 1", dibujar_grafica_planeacion_2024_2026,
                                    d2426, "total_consultas_meta", "total_consultas",
                                    "Consultas totales")
            ft_planeacion(s, lay, "tabla_1", tabla_consultas,
                          w1=4.60, w2=1.35, w3=1.35, w4=1.40,
                          size_header=11, size_body=10, h_fila=0.38)
            _set_texto_placeholder(s, lay, "fecha", f"Del 01 de enero al {fecha_portada}")

    # Diapo 5: serie mensual de consultas ----------------------------------
    sm = historicos.copy()
    sm["fecha"] = pd.to_datetime(sm["fecha"]).dt.tz_localize(None)
    sm["fecha"] = sm["fecha"].apply(lambda d: pd.Timestamp(_floor_month(d)))
    sm = sm[sm["fecha"].notna() & sm["consulta_total"].notna()]
    serie_mensual_consultas = (sm.groupby("fecha", as_index=False)
                               .agg(consultas_totales=("consulta_total", "sum"))
                               .sort_values("fecha"))

    subtitulo_5 = f"Agosto 2022 – {MESES_ES[fecha_fin_graf.month - 1].capitalize()} {fecha_fin_graf.year}"

    lay = _buscar_layout(prs, "Una grafica")
    s = prs.slides.add_slide(lay)
    _set_texto_placeholder(s, lay, "Título 1", "Consultas totales por mes")
    _set_texto_placeholder(s, lay, "fecha", subtitulo_5)
    _dibujar_en_placeholder(s, lay, "ft", dibujar_grafica_consultas_periodos,
                            serie_mensual_consultas, fecha_inicio="2022-08-01",
                            fecha_fin=str(fecha_fin_graf))

    # Diapo 6: serie mensual de procedimientos quirúrgicos -----------------
    if hay_indicador_2026(datos_consulta_funcion, "qx"):
        smq = historicos.copy()
        smq["fecha"] = pd.to_datetime(smq["fecha"]).dt.tz_localize(None)
        smq["fecha"] = smq["fecha"].apply(lambda d: pd.Timestamp(_floor_month(d)))
        smq = smq[smq["fecha"].notna() & smq["consulta_total"].notna()]
        serie_mensual_pq = (smq.groupby("fecha", as_index=False)
                            .agg(consultas_totales=("procedimientos_qx", "sum"))
                            .sort_values("fecha"))

        subtitulo_6 = f"Agosto 2022 – {MESES_ES[fecha_fin_graf.month - 1].capitalize()} {fecha_fin_graf.year}"

        lay = _buscar_layout(prs, "Una grafica")
        s = prs.slides.add_slide(lay)
        _set_texto_placeholder(s, lay, "Título 1", "Procedimientos quirúrgicos por mes")
        _set_texto_placeholder(s, lay, "fecha", subtitulo_6)
        _dibujar_en_placeholder(s, lay, "ft", dibujar_grafica_consultas_periodos,
                                serie_mensual_pq, fecha_inicio="2022-08-01",
                                fecha_fin=str(fecha_fin_graf))

    return prs
